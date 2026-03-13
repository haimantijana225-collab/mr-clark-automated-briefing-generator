from pptx import Presentation
from pptx.util import Inches, Pt


def generate_ppt(slides, chart_path=None):

    prs = Presentation()

    for i, slide_points in enumerate(slides):

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = f"Key Insight {i+1}"

        tf = body.text_frame
        tf.clear()

        for j, point in enumerate(slide_points):

            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = point
            p.level = 0
            p.font.size = Pt(22)

        # Chart only on first slide
        if chart_path and i == 0:

            slide.shapes.add_picture(
                chart_path,
                Inches(6.5),
                Inches(2.5),
                height=Inches(3)
            )

    output = "generated_presentation.pptx"

    prs.save(output)

    return output